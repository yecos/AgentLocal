"""
=============================================================
AGENTE v14.7 - Herramientas de Creacion de Documentos
=============================================================
Genera documentos profesionales desde el agente:
- PDF (fpdf2 / reportlab)
- DOCX (python-docx)
- XLSX (openpyxl)
- Graficos y visualizaciones (matplotlib)

Cada herramienta intenta primero la libreria Python,
y si no esta instalada, da instrucciones claras.
=============================================================
"""

import os
import json
import logging
from datetime import datetime

from config import REPOS_DIR, logger
from utils.security import validate_path


# ============================================================
# CREACION DE PDF
# ============================================================

def crear_pdf(ruta: str, titulo: str = "", contenido: str = "") -> str:
    """Crea un documento PDF con texto formateado. Soporta titulos, parrafos y formato basico.

    Args:
        ruta: Ruta donde guardar el PDF
        titulo: Titulo del documento
        contenido: Contenido del documento (soporta saltos de linea)
    """
    validation = validate_path(ruta)
    if validation != ruta:
        return validation

    # Intentar con fpdf2
    result = _crear_pdf_fpdf(ruta, titulo, contenido)
    if result is not None:
        return result

    # Intentar con reportlab
    result = _crear_pdf_reportlab(ruta, titulo, contenido)
    if result is not None:
        return result

    return ("ERROR: No se pudo crear el PDF. Instala:\n"
            "  pip install fpdf2   (recomendado, mas simple)\n"
            "  pip install reportlab (alternativa, mas potente)")


def _crear_pdf_fpdf(ruta, titulo, contenido):
    """Crea PDF con fpdf2."""
    try:
        from fpdf import FPDF

        pdf = FPDF()
        pdf.add_page()
        pdf.set_auto_page_break(auto=True, margin=15)

        # Intentar usar fuente con soporte Unicode
        try:
            pdf.add_font("DejaVu", "", "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", uni=True)
            pdf.add_font("DejaVu", "B", "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", uni=True)
            font = "DejaVu"
        except Exception:
            font = "Helvetica"

        # Titulo
        if titulo:
            pdf.set_font(font, "B", 16)
            pdf.cell(0, 10, titulo, new_x="LMARGIN", new_y="NEXT", align="C")
            pdf.ln(5)

        # Fecha
        pdf.set_font(font, "", 10)
        pdf.cell(0, 6, f"Generado: {datetime.now().strftime('%Y-%m-%d %H:%M')}", new_x="LMARGIN", new_y="NEXT", align="R")
        pdf.ln(5)

        # Contenido
        if contenido:
            pdf.set_font(font, "", 11)
            for line in contenido.split("\n"):
                # Detectar headers
                if line.startswith("# "):
                    pdf.set_font(font, "B", 14)
                    pdf.cell(0, 8, line[2:], new_x="LMARGIN", new_y="NEXT")
                    pdf.set_font(font, "", 11)
                elif line.startswith("## "):
                    pdf.set_font(font, "B", 12)
                    pdf.cell(0, 7, line[3:], new_x="LMARGIN", new_y="NEXT")
                    pdf.set_font(font, "", 11)
                elif line.startswith("- ") or line.startswith("* "):
                    pdf.cell(5)
                    pdf.cell(0, 6, line, new_x="LMARGIN", new_y="NEXT")
                elif line.strip() == "":
                    pdf.ln(3)
                else:
                    pdf.multi_cell(0, 6, line)

        # Guardar
        dir_name = os.path.dirname(ruta)
        if dir_name:
            os.makedirs(dir_name, exist_ok=True)

        pdf.output(ruta)
        return f"PDF creado: {ruta} ({os.path.getsize(ruta):,} bytes)"

    except ImportError:
        return None
    except Exception as e:
        logger.debug(f"fpdf2 fallo: {e}")
        return f"ERROR creando PDF: {e}"


def _crear_pdf_reportlab(ruta, titulo, contenido):
    """Crea PDF con reportlab (alternativa)."""
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.units import cm

        doc = SimpleDocTemplate(ruta, pagesize=A4,
                                leftMargin=2*cm, rightMargin=2*cm,
                                topMargin=2*cm, bottomMargin=2*cm)

        styles = getSampleStyleSheet()
        elements = []

        if titulo:
            elements.append(Paragraph(titulo, styles['Title']))
            elements.append(Spacer(1, 0.5*cm))

        elements.append(Paragraph(
            f"Generado: {datetime.now().strftime('%Y-%m-%d %H:%M')}",
            styles['Normal']
        ))
        elements.append(Spacer(1, 0.5*cm))

        if contenido:
            for line in contenido.split("\n"):
                if line.startswith("# "):
                    elements.append(Paragraph(line[2:], styles['Heading1']))
                elif line.startswith("## "):
                    elements.append(Paragraph(line[3:], styles['Heading2']))
                elif line.strip():
                    elements.append(Paragraph(line, styles['Normal']))
                else:
                    elements.append(Spacer(1, 0.3*cm))

        dir_name = os.path.dirname(ruta)
        if dir_name:
            os.makedirs(dir_name, exist_ok=True)

        doc.build(elements)
        return f"PDF creado: {ruta} ({os.path.getsize(ruta):,} bytes)"

    except ImportError:
        return None
    except Exception as e:
        logger.debug(f"reportlab fallo: {e}")
        return f"ERROR creando PDF con reportlab: {e}"


# ============================================================
# CREACION DE DOCX (Word)
# ============================================================

def crear_docx(ruta: str, titulo: str = "", contenido: str = "",
               tabla: str = "", imagen: str = "") -> str:
    """Crea un documento Word (.docx) con formato profesional.
    Soporta titulos, headers, listas, tablas e imagenes.

    Args:
        ruta: Ruta donde guardar el .docx
        titulo: Titulo del documento
        contenido: Contenido del documento (formato Markdown basico)
        tabla: Tabla en formato JSON: {"headers": ["col1", ...], "rows": [[val, ...], ...]}
        imagen: Ruta de imagen para insertar en el documento
    """
    validation = validate_path(ruta)
    if validation != ruta:
        return validation

    try:
        from docx import Document
        from docx.shared import Pt, Inches, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        doc = Document()

        # Titulo
        if titulo:
            heading = doc.add_heading(titulo, level=0)
            heading.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # Fecha
        date_para = doc.add_paragraph()
        date_para.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        date_run = date_para.add_run(f"Generado: {datetime.now().strftime('%Y-%m-%d %H:%M')}")
        date_run.font.size = Pt(9)
        date_run.font.color.rgb = RGBColor(128, 128, 128)

        # Contenido
        if contenido:
            for line in contenido.split("\n"):
                if line.startswith("# "):
                    doc.add_heading(line[2:], level=1)
                elif line.startswith("## "):
                    doc.add_heading(line[3:], level=2)
                elif line.startswith("### "):
                    doc.add_heading(line[4:], level=3)
                elif line.startswith("#### "):
                    doc.add_heading(line[5:], level=4)
                elif line.startswith("- ") or line.startswith("* "):
                    doc.add_paragraph(line[2:], style='List Bullet')
                elif line.startswith("1. ") or line.startswith("1) "):
                    doc.add_paragraph(line[3:], style='List Number')
                elif line.startswith("> "):
                    # Blockquote - estilo italic con borde
                    p = doc.add_paragraph()
                    p.paragraph_format.left_indent = Inches(0.5)
                    run = p.add_run(line[2:])
                    run.font.italic = True
                    run.font.color.rgb = RGBColor(100, 100, 100)
                elif line.startswith("---"):
                    # Separador horizontal
                    p = doc.add_paragraph()
                    p.paragraph_format.space_before = Pt(6)
                    p.paragraph_format.space_after = Pt(6)
                    run = p.add_run("─" * 60)
                    run.font.color.rgb = RGBColor(200, 200, 200)
                elif line.startswith("**") and line.endswith("**"):
                    # Bold text
                    p = doc.add_paragraph()
                    run = p.add_run(line.strip('*'))
                    run.font.bold = True
                elif line.strip() == "":
                    pass
                else:
                    # Parrafo normal - soporte bold inline con **texto**
                    _add_rich_paragraph(doc, line)

        # Tabla
        if tabla:
            _add_table_to_docx(doc, tabla)

        # Imagen
        if imagen:
            _add_image_to_docx(doc, imagen)

        # Guardar
        dir_name = os.path.dirname(ruta)
        if dir_name:
            os.makedirs(dir_name, exist_ok=True)

        doc.save(ruta)
        return f"DOCX creado: {ruta} ({os.path.getsize(ruta):,} bytes)"

    except ImportError:
        return ("ERROR: No se pudo crear el DOCX. Instala:\n"
                "  pip install python-docx")
    except Exception as e:
        return f"ERROR creando DOCX: {e}"


# ============================================================
# CREACION DE XLSX (Excel)
# ============================================================

def crear_xlsx(ruta: str, datos: str = "", hoja: str = "Hoja1",
               formulas: str = "", grafico_embebido: str = "") -> str:
    """Crea un archivo Excel (.xlsx) profesional con datos, formulas y estilos.
    Los datos se pueden pasar como CSV o JSON.

    Args:
        ruta: Ruta donde guardar el .xlsx
        datos: Datos en formato CSV o JSON
        hoja: Nombre de la hoja (default: Hoja1)
        formulas: Formulas en JSON: [{"celda": "C1", "formula": "=SUM(A1:B1)"}, ...]
        grafico_embebido: Configuracion de grafico JSON: {"tipo": "bar", "rango": "A1:C5", "titulo": "Ventas"}
    """
    validation = validate_path(ruta)
    if validation != ruta:
        return validation

    try:
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = hoja[:31]  # Max 31 chars

        # Estilos
        header_font = Font(bold=True, size=11)
        header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
        header_font_white = Font(bold=True, size=11, color="FFFFFF")
        thin_border = Border(
            left=Side(style='thin'), right=Side(style='thin'),
            top=Side(style='thin'), bottom=Side(style='thin')
        )

        # Parsear datos
        rows = _parse_datos(datos)

        if rows:
            # Primera fila como headers con estilo
            headers = rows[0]
            for col, header in enumerate(headers, 1):
                cell = ws.cell(row=1, column=col, value=header)
                cell.font = header_font_white
                cell.fill = header_fill
                cell.alignment = Alignment(horizontal='center')
                cell.border = thin_border

            # Datos
            for row_idx, row in enumerate(rows[1:], 2):
                for col_idx, value in enumerate(row, 1):
                    cell = ws.cell(row=row_idx, column=col_idx, value=value)
                    cell.border = thin_border
                    # Auto-detectar tipo numerico
                    if isinstance(value, str):
                        try:
                            if '.' in value:
                                cell.value = float(value)
                            else:
                                cell.value = int(value)
                        except ValueError:
                            pass

            # Auto-ajustar ancho de columnas
            for col in range(1, len(headers) + 1):
                max_length = len(str(headers[col - 1]))
                for row in ws.iter_rows(min_row=2, max_row=min(len(rows), 100), min_col=col, max_col=col):
                    for cell in row:
                        if cell.value:
                            max_length = max(max_length, len(str(cell.value)))
                ws.column_dimensions[openpyxl.utils.get_column_letter(col)].width = min(max_length + 2, 50)

        # Guardar
        dir_name = os.path.dirname(ruta)
        if dir_name:
            os.makedirs(dir_name, exist_ok=True)

        # Formulas
        if formulas:
            _add_formulas_to_xlsx(ws, formulas)

        # Grafico embebido
        if grafico_embebido:
            _add_chart_to_xlsx(ws, grafico_embebido, len(rows))

        # Formato condicional basico - alternar colores de fila
        alt_fill = PatternFill(start_color="D9E2F3", end_color="D9E2F3", fill_type="solid")
        for row_idx in range(2, ws.max_row + 1):
            if row_idx % 2 == 0:
                for col_idx in range(1, ws.max_column + 1):
                    cell = ws.cell(row=row_idx, column=col_idx)
                    if not cell.fill or cell.fill.start_color.rgb == '00000000':
                        cell.fill = alt_fill

        # Freeze panes (header row)
        ws.freeze_panes = 'A2'

        # Auto-filtro
        if ws.max_row > 1 and ws.max_column > 0:
            ws.auto_filter.ref = f"A1:{openpyxl.utils.get_column_letter(ws.max_column)}{ws.max_row}"

        wb.save(ruta)
        return f"XLSX creado: {ruta} ({os.path.getsize(ruta):,} bytes, {len(rows)} filas, {ws.max_column} columnas)"

    except ImportError:
        return ("ERROR: No se pudo crear el XLSX. Instala:\n"
                "  pip install openpyxl")
    except Exception as e:
        return f"ERROR creando XLSX: {e}"


def _parse_datos(datos):
    """Parsea datos en formato CSV o JSON a lista de listas."""
    import json
    import csv
    import io

    if not datos:
        return []

    datos = datos.strip()

    # Intentar JSON
    if datos.startswith('[') or datos.startswith('{'):
        try:
            parsed = json.loads(datos)
            if isinstance(parsed, list):
                if isinstance(parsed[0], dict):
                    headers = list(parsed[0].keys())
                    rows = [headers]
                    for item in parsed:
                        rows.append([str(item.get(h, "")) for h in headers])
                    return rows
                elif isinstance(parsed[0], list):
                    return parsed
            return [["data"], [str(parsed)]]
        except (json.JSONDecodeError, IndexError):
            pass

    # Fallback: CSV
    reader = csv.reader(io.StringIO(datos))
    return [row for row in reader]


# ============================================================
# FUNCIONES AUXILIARES PARA DOCX
# ============================================================

def _add_rich_paragraph(doc, text: str):
    """Agrega un parrafo con formato rich (bold inline **texto**)."""
    from docx.shared import Pt, RGBColor

    p = doc.add_paragraph()
    parts = text.split("**")
    for i, part in enumerate(parts):
        if not part:
            continue
        run = p.add_run(part)
        if i % 2 == 1:  # Odd indices are between ** pairs = bold
            run.font.bold = True


def _add_table_to_docx(doc, tabla_json: str):
    """Agrega una tabla formateada al documento DOCX."""
    try:
        from docx.shared import Pt, Inches, RGBColor
        from docx.enum.table import WD_TABLE_ALIGNMENT
        from docx.oxml.ns import qn

        data = json.loads(tabla_json)
        headers = data.get("headers", data.get("columnas", []))
        rows = data.get("rows", data.get("filas", []))

        if not headers and not rows:
            return

        # Crear tabla
        num_cols = len(headers) if headers else max(len(r) for r in rows)
        num_rows = len(rows) + (1 if headers else 0)

        table = doc.add_table(rows=num_rows, cols=num_cols)
        table.alignment = WD_TABLE_ALIGNMENT.CENTER
        table.style = 'Light Grid Accent 1'

        # Headers
        if headers:
            for i, header in enumerate(headers):
                cell = table.rows[0].cells[i]
                cell.text = str(header)
                for paragraph in cell.paragraphs:
                    paragraph.alignment = 1  # Center
                    for run in paragraph.runs:
                        run.font.bold = True
                        run.font.size = Pt(10)

        # Data rows
        for row_idx, row_data in enumerate(rows):
            table_row = table.rows[row_idx + (1 if headers else 0)]
            for col_idx, value in enumerate(row_data):
                if col_idx < num_cols:
                    cell = table_row.cells[col_idx]
                    cell.text = str(value)
                    for paragraph in cell.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(10)

    except json.JSONDecodeError:
        doc.add_paragraph(f"[Tabla: formato JSON invalido]")
    except Exception as e:
        doc.add_paragraph(f"[Error creando tabla: {e}]")


def _add_image_to_docx(doc, imagen_path: str):
    """Agrega una imagen al documento DOCX."""
    try:
        from docx.shared import Inches

        if os.path.exists(imagen_path):
            doc.add_picture(imagen_path, width=Inches(5.0))
            # Centrar imagen
            last_paragraph = doc.paragraphs[-1]
            last_paragraph.alignment = 1  # Center
        else:
            doc.add_paragraph(f"[Imagen no encontrada: {imagen_path}]")
    except Exception as e:
        doc.add_paragraph(f"[Error insertando imagen: {e}]")


# ============================================================
# FUNCIONES AUXILIARES PARA XLSX
# ============================================================

def _add_formulas_to_xlsx(ws, formulas_json: str):
    """Agrega formulas a celdas del Excel."""
    try:
        formulas = json.loads(formulas_json)
        if isinstance(formulas, list):
            for f in formulas:
                celda = f.get("celda", f.get("cell", ""))
                formula = f.get("formula", f.get("expresion", ""))
                if celda and formula:
                    try:
                        ws[celda] = formula
                    except Exception:
                        pass
        elif isinstance(formulas, dict):
            # Formato: {"C1": "=SUM(A1:B1)", "C2": "=AVERAGE(A1:B1)"}
            for celda, formula in formulas.items():
                try:
                    ws[celda] = str(formula)
                except Exception:
                    pass
    except json.JSONDecodeError:
        pass
    except Exception:
        pass


def _add_chart_to_xlsx(ws, chart_json: str, num_rows: int):
    """Agrega un grafico embebido al Excel."""
    try:
        from openpyxl.chart import BarChart, LineChart, PieChart, Reference

        config = json.loads(chart_json)
        chart_type = config.get("tipo", config.get("type", "bar"))
        data_range = config.get("rango", config.get("range", ""))
        title = config.get("titulo", config.get("title", ""))

        if not data_range:
            # Auto-rango si no se especifica
            from openpyxl.utils import get_column_letter
            max_col = ws.max_column
            max_row = ws.max_row
            data_range = f"A1:{get_column_letter(max_col)}{max_row}"

        # Crear tipo de grafico
        if chart_type in ("pie", "torta"):
            chart = PieChart()
        elif chart_type in ("line", "linea"):
            chart = LineChart()
        else:
            chart = BarChart()

        if title:
            chart.title = title

        # Extraer datos del rango
        # Formato: "A1:C5" -> min_col, min_row, max_col, max_row
        import re
        match = re.match(r'([A-Z]+)(\d+):([A-Z]+)(\d+)', data_range)
        if match:
            from openpyxl.utils import column_index_from_string
            min_col = column_index_from_string(match.group(1))
            min_row = int(match.group(2))
            max_col = column_index_from_string(match.group(3))
            max_row = int(match.group(4))

            data = Reference(ws, min_col=min_col, min_row=min_row,
                           max_col=max_col, max_row=max_row)
            cats = Reference(ws, min_col=min_col, min_row=min_row + 1,
                           max_row=max_row)

            chart.add_data(data, titles_from_data=True)
            chart.set_categories(cats)

            ws.add_chart(chart, f"{get_column_letter(max_col + 2)}1")

    except ImportError:
        pass
    except Exception as e:
        logger.debug(f"Error agregando grafico a XLSX: {e}")


# ============================================================
# CREACION DE GRAFICOS (matplotlib)
# ============================================================

def crear_grafico(ruta: str, tipo: str = "bar", datos: str = "",
                  titulo: str = "", xlabel: str = "", ylabel: str = "") -> str:
    """Crea un grafico o visualizacion y lo guarda como imagen PNG.

    Args:
        ruta: Ruta donde guardar la imagen PNG
        tipo: Tipo de grafico: bar, line, pie, scatter, histogram, area
        datos: Datos en formato CSV (columnas: etiqueta,valor) o JSON
        titulo: Titulo del grafico
        xlabel: Etiqueta del eje X
        ylabel: Etiqueta del eje Y
    """
    validation = validate_path(ruta)
    if validation != ruta:
        return validation

    try:
        import matplotlib
        matplotlib.use('Agg')  # Backend sin GUI
        import matplotlib.pyplot as plt
        import matplotlib.font_manager as fm

        # Configurar fuentes para espanol
        fm.fontManager.addfont('/usr/share/fonts/truetype/chinese/NotoSansSC[wght].ttf')
        fm.fontManager.addfont('/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf')
        plt.rcParams['font.sans-serif'] = ['Noto Sans SC', 'DejaVu Sans']
        plt.rcParams['axes.unicode_minus'] = False

        # Parsear datos
        labels, values = _parse_datos_grafico(datos)

        if not labels or not values:
            return "ERROR: No se pudieron parsear los datos. Usa formato: etiqueta,valor por linea."

        # Crear grafico segun tipo
        fig, ax = plt.subplots(figsize=(10, 6))

        if tipo == "bar":
            ax.bar(labels, values, color='#4472C4')
            if xlabel:
                ax.set_xlabel(xlabel)
            if ylabel:
                ax.set_ylabel(ylabel)
            plt.xticks(rotation=45, ha='right')

        elif tipo == "line":
            ax.plot(labels, values, marker='o', color='#4472C4', linewidth=2)
            ax.fill_between(range(len(labels)), values, alpha=0.1, color='#4472C4')
            if xlabel:
                ax.set_xlabel(xlabel)
            if ylabel:
                ax.set_ylabel(ylabel)
            plt.xticks(rotation=45, ha='right')

        elif tipo == "pie":
            colors = ['#4472C4', '#ED7D31', '#A5A5A5', '#FFC000', '#5B9BD5', '#70AD47', '#264478', '#9B59B6']
            ax.pie(values, labels=labels, autopct='%1.1f%%', colors=colors[:len(labels)])
            ax.axis('equal')

        elif tipo == "scatter":
            if len(values) > 1:
                ax.scatter(labels[:len(values)], values, color='#4472C4', s=100, alpha=0.7)
            else:
                return "ERROR: Scatter necesita al menos 2 series de datos."
            if xlabel:
                ax.set_xlabel(xlabel)
            if ylabel:
                ax.set_ylabel(ylabel)

        elif tipo == "histogram":
            ax.hist(values, bins=min(20, len(values)), color='#4472C4', edgecolor='white')
            if xlabel:
                ax.set_xlabel(xlabel)
            if ylabel:
                ax.set_ylabel(ylabel)

        elif tipo == "area":
            ax.fill_between(range(len(labels)), values, alpha=0.3, color='#4472C4')
            ax.plot(range(len(labels)), values, color='#4472C4', linewidth=2)
            ax.set_xticks(range(len(labels)))
            ax.set_xticklabels(labels, rotation=45, ha='right')
            if xlabel:
                ax.set_xlabel(xlabel)
            if ylabel:
                ax.set_ylabel(ylabel)

        else:
            plt.close()
            return f"ERROR: Tipo de grafico '{tipo}' no soportado. Usar: bar, line, pie, scatter, histogram, area"

        if titulo:
            ax.set_title(titulo, fontsize=14, fontweight='bold')

        plt.tight_layout()

        # Guardar
        dir_name = os.path.dirname(ruta)
        if dir_name:
            os.makedirs(dir_name, exist_ok=True)

        plt.savefig(ruta, dpi=150, bbox_inches='tight')
        plt.close()

        return f"Grafico creado: {ruta} ({os.path.getsize(ruta):,} bytes)"

    except ImportError:
        return ("ERROR: No se pudo crear el grafico. Instala:\n"
                "  pip install matplotlib")
    except Exception as e:
        return f"ERROR creando grafico: {e}"


def _parse_datos_grafico(datos):
    """Parsea datos para graficos. Retorna (labels, values)."""
    import json

    if not datos:
        return [], []

    datos = datos.strip()

    # Intentar JSON
    if datos.startswith('{') or datos.startswith('['):
        try:
            parsed = json.loads(datos)
            if isinstance(parsed, dict):
                return list(parsed.keys()), list(parsed.values())
            elif isinstance(parsed, list):
                if isinstance(parsed[0], dict):
                    labels = [str(d.get('label', d.get('etiqueta', d.get('name', '')))) for d in parsed]
                    values = [float(d.get('value', d.get('valor', d.get('count', 0)))) for d in parsed]
                    return labels, values
                elif isinstance(parsed[0], list):
                    labels = [str(r[0]) for r in parsed]
                    values = [float(r[1]) if len(r) > 1 else 0 for r in parsed]
                    return labels, values
        except (json.JSONDecodeError, IndexError, ValueError):
            pass

    # Fallback: CSV (etiqueta,valor por linea)
    labels = []
    values = []
    for line in datos.split("\n"):
        line = line.strip()
        if not line or line.startswith("#"):
            continue
        # Separar por coma
        parts = line.rsplit(",", 1)
        if len(parts) == 2:
            labels.append(parts[0].strip())
            try:
                values.append(float(parts[1].strip()))
            except ValueError:
                values.append(0)
        else:
            labels.append(line)
            values.append(0)

    return labels, values


# ============================================================
# CREACION DE PPTX (PowerPoint)
# ============================================================

def crear_pptx(ruta: str, titulo: str = "", diapositivas: str = "[]",
               autor: str = "") -> str:
    """Crea una presentacion PowerPoint (.pptx) con titulo y diapositivas.

    Args:
        ruta: Ruta donde guardar el .pptx
        titulo: Titulo de la presentacion
        diapositivas: Lista JSON de diapositivas: [{"titulo": "...", "contenido": "...", "notas": "..."}]
        autor: Autor de la presentacion
    """
    validation = validate_path(ruta)
    if validation != ruta:
        return validation

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Diapositiva de titulo
        if titulo:
            slide_layout = prs.slide_layouts[0]  # Title slide
            slide = prs.slides.add_slide(slide_layout)
            title_shape = slide.shapes.title
            subtitle_shape = slide.placeholders[1]

            title_shape.text = titulo
            for paragraph in title_shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(40)
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(0x44, 0x72, 0xC4)

            if autor:
                subtitle_shape.text = autor
                for paragraph in subtitle_shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(20)
                        run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

        # Parsear diapositivas
        try:
            slides_data = json.loads(diapositivas) if diapositivas else []
        except json.JSONDecodeError:
            # Si no es JSON, tratar como texto separado por "---"
            slides_data = [{"titulo": "", "contenido": s.strip()}
                           for s in diapositivas.split("---") if s.strip()]

        for slide_data in slides_data:
            slide_title = slide_data.get("titulo", slide_data.get("title", ""))
            slide_content = slide_data.get("contenido", slide_data.get("content", ""))
            slide_notes = slide_data.get("notas", slide_data.get("notes", ""))

            if slide_title and not slide_content:
                # Diapositiva de seccion
                slide_layout = prs.slide_layouts[2]  # Section header
                slide = prs.slides.add_slide(slide_layout)
                title_shape = slide.shapes.title
                title_shape.text = slide_title
                for paragraph in title_shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(32)
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(0x44, 0x72, 0xC4)
            else:
                # Diapositiva de titulo + contenido
                slide_layout = prs.slide_layouts[1]  # Title and content
                slide = prs.slides.add_slide(slide_layout)

                if slide_title:
                    title_shape = slide.shapes.title
                    title_shape.text = slide_title
                    for paragraph in title_shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(28)
                            run.font.bold = True
                            run.font.color.rgb = RGBColor(0x44, 0x72, 0xC4)

                if slide_content:
                    body_shape = slide.placeholders[1]
                    tf = body_shape.text_frame

                    # Parsear contenido: soporta bullets y texto
                    lines = slide_content.split("\n")
                    for i, line in enumerate(lines):
                        line = line.strip()
                        if not line:
                            continue

                        if i == 0:
                            p = tf.paragraphs[0]
                        else:
                            p = tf.add_paragraph()

                        # Detectar bullets
                        if line.startswith("- ") or line.startswith("* "):
                            p.text = line[2:]
                            p.level = 0
                            p.space_before = Pt(4)
                            p.space_after = Pt(4)
                        elif line.startswith("  - ") or line.startswith("  * "):
                            p.text = line.strip()[2:]
                            p.level = 1
                            p.space_before = Pt(2)
                            p.space_after = Pt(2)
                        elif line.startswith("# "):
                            p.text = line[2:]
                            for run in p.runs:
                                run.font.size = Pt(20)
                                run.font.bold = True
                        elif line.startswith("## "):
                            p.text = line[3:]
                            for run in p.runs:
                                run.font.size = Pt(16)
                                run.font.bold = True
                        else:
                            p.text = line

                        for run in p.runs:
                            run.font.size = Pt(16)

            # Notas del orador
            if slide_notes:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = slide_notes

        # Guardar
        dir_name = os.path.dirname(ruta)
        if dir_name:
            os.makedirs(dir_name, exist_ok=True)

        prs.save(ruta)
        n_slides = len(prs.slides)
        size_kb = os.path.getsize(ruta) / 1024
        return f"PPTX creado: {ruta} ({size_kb:.0f} KB, {n_slides} diapositivas)"

    except ImportError:
        return ("ERROR: python-pptx no instalado. Instala:\n"
                "  pip install python-pptx")
    except Exception as e:
        return f"ERROR creando PPTX: {e}"
