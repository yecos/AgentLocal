"""
=============================================================
AGENTE v14.7 - Herramientas de Documentos (Super Agente)
=============================================================
Lectura y extraccion de texto de documentos de todo tipo:
- PDF (PyMuPDF / pdfplumber / comando)
- DOCX (python-docx / comando)
- XLSX (openpyxl / comando)
- PPTX (python-pptx / comando)
- CSV estructurado
- RTF, ePub

Cada herramienta intenta primero la libreria Python,
y si no esta instalada, hace fallback a comandos del sistema.
=============================================================
"""

import os
import csv
import io
import logging
from config import REPOS_DIR, MAX_FILE_READ, MAX_TOOL_OUTPUT, logger
from utils.security import validate_path


# ============================================================
# LECTURA DE PDF
# ============================================================

def leer_pdf(ruta: str, pagina_inicio: int = None, pagina_fin: int = None) -> str:
    """Lee el contenido de un archivo PDF. Extrae texto, tablas e info de paginas.

    Args:
        ruta: Ruta del archivo PDF
        pagina_inicio: Pagina inicial (1-indexed, opcional)
        pagina_fin: Pagina final (inclusive, opcional)
    """
    validation = validate_path(ruta)
    if validation != ruta:
        return validation

    ruta = _resolve_path(ruta)
    if not ruta:
        return f"Archivo no encontrado: {ruta}"

    if not ruta.lower().endswith('.pdf'):
        return "ERROR: El archivo no es un PDF."

    # Intentar con PyMuPDF (fitz)
    result = _leer_pdf_fitz(ruta, pagina_inicio, pagina_fin)
    if result is not None:
        return result

    # Intentar con pdfplumber
    result = _leer_pdf_plumber(ruta, pagina_inicio, pagina_fin)
    if result is not None:
        return result

    # Fallback: comando pdftotext
    result = _leer_pdf_comando(ruta, pagina_inicio, pagina_fin)
    if result is not None:
        return result

    return ("ERROR: No se pudo leer el PDF. Instala una libreria:\n"
            "  pip install PyMuPDF   (recomendado)\n"
            "  pip install pdfplumber (alternativa)\n"
            "  O instala pdftotext: sudo apt install poppler-utils")


def _leer_pdf_fitz(ruta, pagina_inicio, pagina_fin):
    """Lee PDF con PyMuPDF (fitz)."""
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(ruta)

        total_pages = len(doc)
        start = max(0, (pagina_inicio or 1) - 1)
        end = min(total_pages, pagina_fin or total_pages)

        parts = [f"PDF: {os.path.basename(ruta)} ({total_pages} paginas)\n"]

        for i in range(start, end):
            page = doc[i]
            text = page.get_text("text")
            if text.strip():
                parts.append(f"--- Pagina {i + 1} ---")
                parts.append(text)

            # Extraer tablas si hay
            try:
                tables = page.find_tables()
                if tables and tables.tables:
                    for ti, table in enumerate(tables.tables):
                        table_data = table.extract()
                        if table_data:
                            parts.append(f"  Tabla {ti + 1}:")
                            for row in table_data[:20]:
                                parts.append("  | " + " | ".join(str(c or "") for c in row) + " |")
            except Exception:
                pass  # Tabla extraction no critica

        doc.close()

        content = "\n".join(parts)
        if len(content) > MAX_FILE_READ:
            content = content[:MAX_FILE_READ] + f"\n... [truncado, {total_pages} paginas total]"

        return content

    except ImportError:
        return None
    except Exception as e:
        logger.debug(f"PyMuPDF fallo: {e}")
        return None


def _leer_pdf_plumber(ruta, pagina_inicio, pagina_fin):
    """Lee PDF con pdfplumber (alternativa)."""
    try:
        import pdfplumber

        with pdfplumber.open(ruta) as pdf:
            total_pages = len(pdf.pages)
            start = max(0, (pagina_inicio or 1) - 1)
            end = min(total_pages, pagina_fin or total_pages)

            parts = [f"PDF: {os.path.basename(ruta)} ({total_pages} paginas)\n"]

            for i in range(start, end):
                page = pdf.pages[i]
                text = page.extract_text()
                if text and text.strip():
                    parts.append(f"--- Pagina {i + 1} ---")
                    parts.append(text)

                # Extraer tablas
                try:
                    tables = page.extract_tables()
                    for ti, table in enumerate(tables[:3]):
                        parts.append(f"  Tabla {ti + 1}:")
                        for row in table[:20]:
                            parts.append("  | " + " | ".join(str(c or "") for c in row) + " |")
                except Exception:
                    pass

            content = "\n".join(parts)
            if len(content) > MAX_FILE_READ:
                content = content[:MAX_FILE_READ] + f"\n... [truncado, {total_pages} paginas total]"
            return content

    except ImportError:
        return None
    except Exception as e:
        logger.debug(f"pdfplumber fallo: {e}")
        return None


def _leer_pdf_comando(ruta, pagina_inicio, pagina_fin):
    """Fallback: lee PDF con comando pdftotext."""
    import subprocess
    try:
        result = subprocess.run(
            ["pdftotext", "-layout", ruta, "-"],
            capture_output=True, text=True, timeout=30
        )
        if result.returncode == 0 and result.stdout.strip():
            content = f"PDF: {os.path.basename(ruta)}\n\n{result.stdout}"
            if len(content) > MAX_FILE_READ:
                content = content[:MAX_FILE_READ] + "\n... [truncado]"
            return content
    except (FileNotFoundError, subprocess.TimeoutExpired):
        pass
    except Exception as e:
        logger.debug(f"pdftotext fallo: {e}")
    return None


# ============================================================
# LECTURA DE DOCX (Word)
# ============================================================

def leer_docx(ruta: str) -> str:
    """Lee el contenido de un archivo Word (.docx). Extrae texto, tablas y estilos.

    Args:
        ruta: Ruta del archivo .docx
    """
    validation = validate_path(ruta)
    if validation != ruta:
        return validation

    ruta = _resolve_path(ruta)
    if not ruta:
        return f"Archivo no encontrado: {ruta}"

    if not ruta.lower().endswith('.docx'):
        return "ERROR: El archivo no es un .docx (para .doc viejo usar comando antiword)."

    # Intentar con python-docx
    result = _leer_docx_python(ruta)
    if result is not None:
        return result

    # Fallback: comando
    return _leer_docx_comando(ruta)


def _leer_docx_python(ruta):
    """Lee DOCX con python-docx."""
    try:
        from docx import Document

        doc = Document(ruta)
        parts = [f"DOCX: {os.path.basename(ruta)}\n"]

        # Extraer parrafos
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                # Indicar estilo si es heading
                style_name = para.style.name if para.style else ""
                if "Heading" in style_name or "heading" in style_name:
                    level = "".join(c for c in style_name if c.isdigit()) or "1"
                    parts.append(f"{'#' * int(level)} {text}")
                else:
                    parts.append(text)

        # Extraer tablas
        for ti, table in enumerate(doc.tables):
            parts.append(f"\nTabla {ti + 1}:")
            for row in table.rows[:20]:
                cells = [cell.text.strip() for cell in row.cells]
                parts.append("| " + " | ".join(cells) + " |")

        content = "\n".join(parts)
        if len(content) > MAX_FILE_READ:
            content = content[:MAX_FILE_READ] + "\n... [truncado]"
        return content

    except ImportError:
        return None
    except Exception as e:
        logger.debug(f"python-docx fallo: {e}")
        return f"ERROR leyendo DOCX: {e}"


def _leer_docx_comando(ruta):
    """Fallback: lee DOCX con comando (unzip + xml parse)."""
    import subprocess
    try:
        # docx es un ZIP con word/document.xml
        result = subprocess.run(
            ["python3", "-c",
             f"import zipfile,xml.etree.ElementTree as ET;"
             f"z=zipfile.ZipFile('{ruta}');"
             f"tree=ET.parse(z.open('word/document.xml'));"
             f"root=tree.getroot();"
             f"ns={{'w':'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}};"
             f"print('\\n'.join(t.text for t in root.findall('.//w:t',ns) if t.text))"],
            capture_output=True, text=True, timeout=30
        )
        if result.returncode == 0 and result.stdout.strip():
            content = f"DOCX: {os.path.basename(ruta)}\n\n{result.stdout}"
            if len(content) > MAX_FILE_READ:
                content = content[:MAX_FILE_READ] + "\n... [truncado]"
            return content
    except Exception:
        pass

    return ("ERROR: No se pudo leer el DOCX. Instala:\n"
            "  pip install python-docx   (recomendado)")


# ============================================================
# LECTURA DE XLSX (Excel)
# ============================================================

def leer_xlsx(ruta: str, hoja: str = None, max_filas: int = 50) -> str:
    """Lee el contenido de un archivo Excel (.xlsx). Extrae datos de hojas y tablas.

    Args:
        ruta: Ruta del archivo .xlsx
        hoja: Nombre de la hoja (opcional, por defecto la primera)
        max_filas: Maximo de filas a leer por hoja (default 50)
    """
    validation = validate_path(ruta)
    if validation != ruta:
        return validation

    ruta = _resolve_path(ruta)
    if not ruta:
        return f"Archivo no encontrado: {ruta}"

    if not ruta.lower().endswith(('.xlsx', '.xls')):
        return "ERROR: El archivo no es un Excel (.xlsx/.xls)."

    # Intentar con openpyxl
    result = _leer_xlsx_openpyxl(ruta, hoja, max_filas)
    if result is not None:
        return result

    # Fallback: comando
    return _leer_xlsx_comando(ruta, hoja, max_filas)


def _leer_xlsx_openpyxl(ruta, hoja, max_filas):
    """Lee XLSX con openpyxl."""
    try:
        from openpyxl import load_workbook

        wb = load_workbook(ruta, read_only=True, data_only=True)
        parts = [f"XLSX: {os.path.basename(ruta)} ({len(wb.sheetnames)} hojas)\n"]
        parts.append(f"Hojas: {', '.join(wb.sheetnames)}\n")

        # Determinar hojas a leer
        if hoja:
            if hoja in wb.sheetnames:
                sheets = [wb[hoja]]
            else:
                return f"Hoja '{hoja}' no encontrada. Disponibles: {', '.join(wb.sheetnames)}"
        else:
            sheets = [wb[name] for name in wb.sheetnames[:3]]  # Max 3 hojas

        for ws in sheets:
            parts.append(f"=== Hoja: {ws.title} ===")

            row_count = 0
            for row in ws.iter_rows(values_only=True):
                if row_count >= max_filas:
                    parts.append(f"... [truncado, {max_filas} filas maximo]")
                    break
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):  # Skip filas vacias
                    parts.append("| " + " | ".join(cells) + " |")
                    row_count += 1

            parts.append(f"({row_count} filas leidas)\n")

        wb.close()

        content = "\n".join(parts)
        if len(content) > MAX_FILE_READ:
            content = content[:MAX_FILE_READ] + "\n... [truncado]"
        return content

    except ImportError:
        return None
    except Exception as e:
        logger.debug(f"openpyxl fallo: {e}")
        return f"ERROR leyendo XLSX: {e}"


def _leer_xlsx_comando(ruta, hoja, max_filas):
    """Fallback: lee XLSX con comando python + zipfile (basico)."""
    import subprocess
    try:
        python_code = (
            "import zipfile,xml.etree.ElementTree as ET;"
            f"z=zipfile.ZipFile('{ruta}');"
            "tree=ET.parse(z.open('xl/sharedStrings.xml'));"
            "root=tree.getroot();"
            "ns={'s':'http://schemas.openxmlformats.org/spreadsheetml/2006/main'};"
            "strings=[t.text or '' for t in root.findall('.//s:t',ns)];"
            "print('Strings:', len(strings));"
            "print('Sheets:', z.namelist())"
        )
        result = subprocess.run(
            ["python3", "-c", python_code],
            capture_output=True, text=True, timeout=30
        )
        if result.returncode == 0:
            return f"XLSX (basico): {os.path.basename(ruta)}\n{result.stdout}\nInstala openpyxl para lectura completa: pip install openpyxl"
    except Exception:
        pass

    return ("ERROR: No se pudo leer el XLSX. Instala:\n"
            "  pip install openpyxl   (recomendado)")


# ============================================================
# LECTURA DE PPTX (PowerPoint)
# ============================================================

def leer_pptx(ruta: str) -> str:
    """Lee el contenido de una presentacion PowerPoint (.pptx). Extrae texto de diapositivas.

    Args:
        ruta: Ruta del archivo .pptx
    """
    validation = validate_path(ruta)
    if validation != ruta:
        return validation

    ruta = _resolve_path(ruta)
    if not ruta:
        return f"Archivo no encontrado: {ruta}"

    if not ruta.lower().endswith('.pptx'):
        return "ERROR: El archivo no es un .pptx."

    # Intentar con python-pptx
    result = _leer_pptx_python(ruta)
    if result is not None:
        return result

    # Fallback: comando
    return _leer_pptx_comando(ruta)


def _leer_pptx_python(ruta):
    """Lee PPTX con python-pptx."""
    try:
        from pptx import Presentation

        prs = Presentation(ruta)
        parts = [f"PPTX: {os.path.basename(ruta)} ({len(prs.slides)} diapositivas)\n"]

        for i, slide in enumerate(prs.slides):
            parts.append(f"=== Diapositiva {i + 1} ===")

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            parts.append(text)

                # Tablas en diapositivas
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        parts.append("| " + " | ".join(cells) + " |")

            parts.append("")  # Separador entre diapositivas

        content = "\n".join(parts)
        if len(content) > MAX_FILE_READ:
            content = content[:MAX_FILE_READ] + "\n... [truncado]"
        return content

    except ImportError:
        return None
    except Exception as e:
        logger.debug(f"python-pptx fallo: {e}")
        return f"ERROR leyendo PPTX: {e}"


def _leer_pptx_comando(ruta):
    """Fallback: lee PPTX con unzip + xml."""
    import subprocess
    try:
        result = subprocess.run(
            ["python3", "-c",
             f"import zipfile,xml.etree.ElementTree as ET;"
             f"z=zipfile.ZipFile('{ruta}');"
             f"ns={{'a':'http://schemas.openxmlformats.org/drawingml/2006/main'}};"
             f"slides=[n for n in z.namelist() if n.startswith('ppt/slides/slide') and n.endswith('.xml')];"
             f"slides.sort();"
             f"print(f'{{len(slides)}} diapositivas');"
             f"[print(f'Slide {{i+1}}:', '\\n'.join(t.text for t in ET.parse(z.open(s)).getroot().findall('.//a:t',ns) if t.text)) for i,s in enumerate(slates[:20])]"],
            capture_output=True, text=True, timeout=30
        )
        if result.returncode == 0 and result.stdout.strip():
            return f"PPTX (basico): {os.path.basename(ruta)}\n{result.stdout}"
    except Exception:
        pass

    return ("ERROR: No se pudo leer el PPTX. Instala:\n"
            "  pip install python-pptx   (recomendado)")


# ============================================================
# LECTURA DE CSV ESTRUCTURADO
# ============================================================

def leer_csv(ruta: str, separador: str = ",", max_filas: int = 50) -> str:
    """Lee un archivo CSV de forma estructurada. Detecta delimitador y codificacion.

    Args:
        ruta: Ruta del archivo CSV
        separador: Separador de columnas (default: coma)
        max_filas: Maximo de filas a leer (default 50)
    """
    validation = validate_path(ruta)
    if validation != ruta:
        return validation

    ruta = _resolve_path(ruta)
    if not ruta:
        return f"Archivo no encontrado: {ruta}"

    try:
        # Intentar detectar codificacion
        encodings = ['utf-8', 'latin-1', 'cp1252', 'utf-16']
        content = None

        for enc in encodings:
            try:
                with open(ruta, 'r', encoding=enc) as f:
                    content = f.read()
                break
            except (UnicodeDecodeError, UnicodeError):
                continue

        if content is None:
            return "ERROR: No se pudo detectar la codificacion del CSV."

        # Parsear CSV
        reader = csv.reader(io.StringIO(content), delimiter=separador)
        rows = list(reader)

        if not rows:
            return "CSV vacio."

        headers = rows[0] if rows else []
        data_rows = rows[1:max_filas + 1]

        parts = [f"CSV: {os.path.basename(ruta)} ({len(rows) - 1} filas, {len(headers)} columnas)\n"]
        parts.append(f"Columnas: {', '.join(headers)}\n")

        # Formato tabla
        for row in data_rows:
            # Rellenar filas cortas
            while len(row) < len(headers):
                row.append("")
            parts.append("| " + " | ".join(row[:len(headers)]) + " |")

        if len(rows) - 1 > max_filas:
            parts.append(f"\n... [truncado, mostrando {max_filas} de {len(rows) - 1} filas]")

        result = "\n".join(parts)
        if len(result) > MAX_FILE_READ:
            result = result[:MAX_FILE_READ] + "\n... [truncado]"
        return result

    except Exception as e:
        return f"ERROR leyendo CSV: {e}"


# ============================================================
# LECTURA DE ARCHIVOS COMPRIMIDOS
# ============================================================

def leer_archivo_comprimido(ruta: str, archivo_interno: str = None) -> str:
    """Lista o extrae contenido de archivos ZIP, TAR, GZ, RAR.

    Args:
        ruta: Ruta del archivo comprimido
        archivo_interno: Archivo interno a extraer (opcional, si no se especifica lista contenido)
    """
    validation = validate_path(ruta)
    if validation != ruta:
        return validation

    ruta = _resolve_path(ruta)
    if not ruta:
        return f"Archivo no encontrado: {ruta}"

    ext = ruta.lower()

    if ext.endswith('.zip'):
        return _leer_zip(ruta, archivo_interno)
    elif ext.endswith(('.tar', '.tar.gz', '.tgz', '.tar.bz2', '.tbz2')):
        return _leer_tar(ruta, archivo_interno)
    else:
        return "ERROR: Formato no soportado. Usar .zip, .tar, .tar.gz, .tar.bz2"


def _leer_zip(ruta, archivo_interno):
    """Lee contenido de archivo ZIP."""
    try:
        import zipfile
        with zipfile.ZipFile(ruta, 'r') as zf:
            names = zf.namelist()

            if archivo_interno:
                # Extraer archivo interno especifico
                if archivo_interno in names:
                    with zf.open(archivo_interno) as f:
                        content = f.read()
                        try:
                            text = content.decode('utf-8')
                        except UnicodeDecodeError:
                            text = content.decode('latin-1')

                        if len(text) > MAX_FILE_READ:
                            text = text[:MAX_FILE_READ] + "\n... [truncado]"
                        return f"Contenido de {archivo_interno}:\n\n{text}"
                else:
                    return f"Archivo '{archivo_interno}' no encontrado en ZIP. Disponibles:\n" + "\n".join(f"  {n}" for n in names[:30])
            else:
                # Listar contenido
                parts = [f"ZIP: {os.path.basename(ruta)} ({len(names)} archivos)\n"]
                for name in names[:50]:
                    info = zf.getinfo(name)
                    size = info.file_size
                    parts.append(f"  {'[DIR]' if name.endswith('/') else '[FILE]'} {name} ({size:,} bytes)")

                if len(names) > 50:
                    parts.append(f"\n... y {len(names) - 50} archivos mas")

                return "\n".join(parts)

    except Exception as e:
        return f"ERROR leyendo ZIP: {e}"


def _leer_tar(ruta, archivo_interno):
    """Lee contenido de archivo TAR."""
    try:
        import tarfile
        with tarfile.open(ruta, 'r:*') as tf:
            members = tf.getmembers()

            if archivo_interno:
                for member in members:
                    if member.name == archivo_interno:
                        f = tf.extractfile(member)
                        if f:
                            content = f.read()
                            try:
                                text = content.decode('utf-8')
                            except UnicodeDecodeError:
                                text = content.decode('latin-1')
                            if len(text) > MAX_FILE_READ:
                                text = text[:MAX_FILE_READ] + "\n... [truncado]"
                            return f"Contenido de {archivo_interno}:\n\n{text}"
                return f"Archivo '{archivo_interno}' no encontrado en TAR."

            # Listar contenido
            parts = [f"TAR: {os.path.basename(ruta)} ({len(members)} archivos)\n"]
            for m in members[:50]:
                type_str = "[DIR]" if m.isdir() else "[FILE]"
                parts.append(f"  {type_str} {m.name} ({m.size:,} bytes)")

            if len(members) > 50:
                parts.append(f"\n... y {len(members) - 50} archivos mas")

            return "\n".join(parts)

    except Exception as e:
        return f"ERROR leyendo TAR: {e}"


# ============================================================
# LECTURA DE BASES DE DATOS SQLITE
# ============================================================

def consultar_sqlite(ruta: str, consulta: str = None, tabla: str = None, max_filas: int = 50) -> str:
    """Consulta una base de datos SQLite. Puede listar tablas o ejecutar consultas SELECT.

    Args:
        ruta: Ruta del archivo .db o .sqlite
        consulta: Consulta SQL SELECT (opcional)
        tabla: Nombre de tabla para ver contenido (opcional)
        max_filas: Maximo de filas a retornar (default 50)
    """
    validation = validate_path(ruta)
    if validation != ruta:
        return validation

    ruta = _resolve_path(ruta)
    if not ruta:
        return f"Archivo no encontrado: {ruta}"

    try:
        import sqlite3
        conn = sqlite3.connect(ruta)
        cursor = conn.cursor()

        if consulta:
            # Seguridad: solo permitir SELECT
            consulta_upper = consulta.strip().upper()
            if not consulta_upper.startswith("SELECT") and not consulta_upper.startswith("PRAGMA"):
                conn.close()
                return "ERROR: Solo se permiten consultas SELECT o PRAGMA por seguridad."

            # Bloquear comandos peligrosos
            dangerous = ["DROP", "DELETE", "UPDATE", "INSERT", "ALTER", "CREATE", "ATTACH", "DETACH"]
            for word in dangerous:
                if word in consulta_upper.split():
                    conn.close()
                    return f"ERROR: Comando {word} no permitido por seguridad."

            cursor.execute(consulta)
            rows = cursor.fetchall()
            columns = [desc[0] for desc in cursor.description] if cursor.description else []

            parts = [f"Consulta: {consulta}"]
            if columns:
                parts.append("| " + " | ".join(columns) + " |")
            for row in rows[:max_filas]:
                parts.append("| " + " | ".join(str(c) for c in row) + " |")

            if len(rows) > max_filas:
                parts.append(f"\n... [truncado, {len(rows)} filas total]")

            parts.append(f"\n({len(rows)} filas)")
            conn.close()
            return "\n".join(parts)

        elif tabla:
            # Ver contenido de una tabla
            cursor.execute(f'SELECT * FROM "{tabla}" LIMIT {max_filas}')
            rows = cursor.fetchall()
            columns = [desc[0] for desc in cursor.description]

            parts = [f"Tabla: {tabla} ({len(rows)} filas mostradas)\n"]
            parts.append("| " + " | ".join(columns) + " |")
            for row in rows:
                parts.append("| " + " | ".join(str(c) for c in row) + " |")

            conn.close()
            return "\n".join(parts)

        else:
            # Listar tablas
            cursor.execute("SELECT name FROM sqlite_master WHERE type='table' ORDER BY name")
            tables = cursor.fetchall()

            parts = [f"SQLite: {os.path.basename(ruta)} ({len(tables)} tablas)\n"]

            for (table_name,) in tables:
                cursor.execute(f'SELECT COUNT(*) FROM "{table_name}"')
                count = cursor.fetchone()[0]
                parts.append(f"  {table_name} ({count:,} filas)")

                # Mostrar schema
                cursor.execute(f'PRAGMA table_info("{table_name}")')
                cols = cursor.fetchall()
                col_info = ", ".join(f"{c[1]}({c[2]})" for c in cols[:8])
                parts.append(f"    Columnas: {col_info}")

            conn.close()
            return "\n".join(parts)

    except ImportError:
        return "ERROR: sqlite3 no disponible (viene con Python estandar, esto no deberia pasar)."
    except Exception as e:
        return f"ERROR consultando SQLite: {e}"


# ============================================================
# LECTURA DE ePub
# ============================================================

def leer_epub(ruta: str, max_capitulos: int = 10) -> str:
    """Lee el contenido de un libro electronico ePub. Extrae texto de capitulos.

    Args:
        ruta: Ruta del archivo .epub
        max_capitulos: Maximo de capitulos a leer (default 10)
    """
    validation = validate_path(ruta)
    if validation != ruta:
        return validation

    ruta = _resolve_path(ruta)
    if not ruta:
        return f"Archivo no encontrado: {ruta}"

    # Intentar con ebooklib
    try:
        import ebooklib
        from ebooklib import epub
        from html.parser import HTMLParser

        class _TextExtractor(HTMLParser):
            def __init__(self):
                super().__init__()
                self.text = []
                self.skip = False

            def handle_starttag(self, tag, attrs):
                if tag in ('script', 'style'):
                    self.skip = True

            def handle_endtag(self, tag):
                if tag in ('script', 'style'):
                    self.skip = False
                elif tag in ('p', 'div', 'h1', 'h2', 'h3', 'br'):
                    self.text.append('\n')

            def handle_data(self, data):
                if not self.skip:
                    self.text.append(data.strip())

        book = epub.read_epub(ruta)
        parts = [f"ePub: {os.path.basename(ruta)}\n"]

        # Metadata
        title = book.get_metadata('DC', 'title')
        author = book.get_metadata('DC', 'author')
        if title:
            parts.append(f"Titulo: {title[0][0]}")
        if author:
            parts.append(f"Autor: {author[0][0]}")

        # Extraer texto de documentos
        chapters = 0
        for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
            if chapters >= max_capitulos:
                break
            content = item.get_content().decode('utf-8', errors='replace')
            extractor = _TextExtractor()
            extractor.feed(content)
            text = " ".join(extractor.text).strip()
            if text and len(text) > 50:
                parts.append(f"\n--- Capitulo {chapters + 1} ---")
                parts.append(text[:2000])
                chapters += 1

        content = "\n".join(parts)
        if len(content) > MAX_FILE_READ:
            content = content[:MAX_FILE_READ] + "\n... [truncado]"
        return content

    except ImportError:
        # Fallback: tratar como ZIP y buscar HTML
        return _leer_epub_zip(ruta, max_capitulos)
    except Exception as e:
        return f"ERROR leyendo ePub: {e}"


def _leer_epub_zip(ruta, max_capitulos):
    """Fallback: lee ePub como ZIP (formato basico)."""
    try:
        import zipfile
        from html.parser import HTMLParser

        class _SimpleText(HTMLParser):
            def __init__(self):
                super().__init__()
                self.text = []
            def handle_data(self, data):
                self.text.append(data.strip())

        with zipfile.ZipFile(ruta, 'r') as zf:
            html_files = [n for n in zf.namelist() if n.endswith(('.html', '.xhtml'))]

            parts = [f"ePub (basico): {os.path.basename(ruta)} ({len(html_files)} secciones)\n"]

            for i, name in enumerate(html_files[:max_capitulos]):
                content = zf.open(name).read().decode('utf-8', errors='replace')
                extractor = _SimpleText()
                extractor.feed(content)
                text = " ".join(t for t in extractor.text if t)
                if text:
                    parts.append(f"--- Seccion {i + 1} ---")
                    parts.append(text[:1500])

            content = "\n".join(parts)
            if len(content) > MAX_FILE_READ:
                content = content[:MAX_FILE_READ] + "\n... [truncado]"
            return content

    except Exception as e:
        return f"ERROR leyendo ePub: {e}\nInstala: pip install ebooklib"


# ============================================================
# LECTOR UNIVERSAL (auto-detectar formato)
# ============================================================

def leer_documento(ruta: str, **kwargs) -> str:
    """Lee cualquier documento detectando automaticamente el formato.
    Soporta: PDF, DOCX, XLSX, PPTX, CSV, ZIP, TAR, SQLite, ePub, y texto plano.

    Args:
        ruta: Ruta del archivo a leer
        **kwargs: Argumentos adicionales segun el tipo (hoja, consulta, etc.)
    """
    validation = validate_path(ruta)
    if validation != ruta:
        return validation

    ruta_resolved = _resolve_path(ruta)
    if not ruta_resolved:
        return f"Archivo no encontrado: {ruta}"

    ext = ruta_resolved.lower()

    if ext.endswith('.pdf'):
        return leer_pdf(ruta_resolved, **{k: v for k, v in kwargs.items() if k in ('pagina_inicio', 'pagina_fin')})
    elif ext.endswith('.docx'):
        return leer_docx(ruta_resolved)
    elif ext.endswith(('.xlsx', '.xls')):
        return leer_xlsx(ruta_resolved, **{k: v for k, v in kwargs.items() if k in ('hoja', 'max_filas')})
    elif ext.endswith('.pptx'):
        return leer_pptx(ruta_resolved)
    elif ext.endswith('.csv'):
        return leer_csv(ruta_resolved, **{k: v for k, v in kwargs.items() if k in ('separador', 'max_filas')})
    elif ext.endswith(('.zip', '.tar', '.tar.gz', '.tgz', '.tar.bz2', '.tbz2')):
        return leer_archivo_comprimido(
            ruta_resolved,
            **{k: v for k, v in kwargs.items() if k in ('archivo_interno',)}
        )
    elif ext.endswith(('.db', '.sqlite', '.sqlite3')):
        return consultar_sqlite(
            ruta_resolved,
            **{k: v for k, v in kwargs.items() if k in ('consulta', 'tabla', 'max_filas')}
        )
    elif ext.endswith('.epub'):
        return leer_epub(ruta_resolved)
    else:
        # Fallback: intentar como texto
        try:
            with open(ruta_resolved, 'r', encoding='utf-8', errors='replace') as f:
                content = f.read()
            if len(content) > MAX_FILE_READ:
                content = content[:MAX_FILE_READ] + "\n... [truncado]"
            return content
        except Exception as e:
            return f"ERROR: No se pudo leer {ruta_resolved}: {e}"


# ============================================================
# UTILIDADES
# ============================================================

def _resolve_path(ruta):
    """Resuelve la ruta del archivo, buscando en ubicaciones alternativas."""
    if os.path.isabs(ruta) and os.path.exists(ruta):
        return ruta

    # Ruta relativa
    if os.path.exists(ruta):
        return ruta

    # Buscar en REPOS_DIR
    alt = os.path.join(REPOS_DIR, ruta)
    if os.path.exists(alt):
        return alt

    # Buscar en subdirectorios
    try:
        for d in os.listdir(REPOS_DIR):
            alt2 = os.path.join(REPOS_DIR, d, ruta)
            if os.path.exists(alt2):
                return alt2
    except OSError:
        pass

    # Buscar en home
    home = os.path.expanduser("~")
    alt3 = os.path.join(home, ruta)
    if os.path.exists(alt3):
        return alt3

    return None


# ============================================================
# RESUMIR DOCUMENTO
# ============================================================

def resumir_documento(ruta: str, tipo_resumen: str = "ejecutivo", idioma: str = "es") -> str:
    """Resume un documento (PDF, DOCX, TXT, MD). Soporta resumenes ejecutivos, detallados y puntos clave.

    Args:
        ruta: Ruta al documento a resumir
        tipo_resumen: Tipo de resumen: ejecutivo (corto), detallado (largo), puntos_clave (bullet points)
        idioma: Idioma del resumen: es (espanol) o en (ingles)
    """
    if not os.path.exists(ruta):
        return f"ERROR: Archivo no encontrado: {ruta}"

    # Determinar tipo de archivo y leer contenido
    ext = os.path.splitext(ruta)[1].lower()
    content = _read_document_for_summary(ruta, ext)

    if not content or not content.strip():
        return "ERROR: No se pudo extraer texto del documento o esta vacio."

    # Truncar si es muy largo para el LLM
    max_chars = 8000
    if len(content) > max_chars:
        # Chunk: keep beginning and end
        chunk_size = max_chars // 2
        content = content[:chunk_size] + "\n\n[... contenido truncado ...]\n\n" + content[-chunk_size:]

    # Generar resumen con LLM
    summary = _generate_summary(content, tipo_resumen, idioma, os.path.basename(ruta))

    if not summary:
        return "ERROR: No se pudo generar el resumen. El LLM no esta disponible."

    # Add metadata
    word_count = len(content.split())
    header = f"Resumen de: {os.path.basename(ruta)}\n"
    header += f"Tipo: {tipo_resumen} | Palabras originales: ~{word_count} | Idioma: {idioma}\n"
    header += "-" * 50 + "\n\n"

    return header + summary


def _read_document_for_summary(ruta: str, ext: str) -> str:
    """Lee el contenido de un documento para resumirlo."""
    try:
        if ext == ".pdf":
            return leer_pdf(ruta)
        elif ext == ".docx":
            return leer_docx(ruta)
        elif ext == ".xlsx":
            return leer_xlsx(ruta)
        elif ext == ".pptx":
            return leer_pptx(ruta)
        elif ext == ".csv":
            return leer_csv(ruta)
        elif ext == ".epub":
            return leer_epub(ruta)
        elif ext in (".txt", ".md", ".markdown", ".rst", ".log"):
            validation = validate_path(ruta)
            if validation != ruta:
                return validation
            with open(ruta, "r", encoding="utf-8", errors="replace") as f:
                return f.read()
        elif ext == ".json":
            validation = validate_path(ruta)
            if validation != ruta:
                return validation
            with open(ruta, "r", encoding="utf-8", errors="replace") as f:
                import json
                data = json.load(f)
                return json.dumps(data, indent=2, ensure_ascii=False)[:8000]
        else:
            # Try reading as text
            validation = validate_path(ruta)
            if validation != ruta:
                return validation
            with open(ruta, "r", encoding="utf-8", errors="replace") as f:
                return f.read()
    except Exception as e:
        return f"ERROR leyendo documento: {e}"


def _generate_summary(content: str, tipo_resumen: str, idioma: str, filename: str) -> str:
    """Usa el LLM para generar un resumen del contenido."""
    try:
        from llm import ollama

        lang_map = {"es": "espanol", "en": "ingles"}
        lang_name = lang_map.get(idioma, "espanol")

        type_prompts = {
            "ejecutivo": (
                "Genera un resumen ejecutivo conciso (maximo 150 palabras). "
                "Captura las ideas principales, conclusiones y datos clave. "
                "No incluya detalles menores."
            ),
            "detallado": (
                "Genera un resumen detallado (maximo 500 palabras). "
                "Incluye ideas principales, argumentos, datos de soporte, "
                "conclusiones y detalles relevantes. Estructura el resumen en secciones."
            ),
            "puntos_clave": (
                "Extrae los puntos clave del documento en formato de bullet points. "
                "Maximo 10 puntos. Cada punto debe ser conciso pero informativo. "
                "Usa el formato: - Punto clave"
            ),
        }

        summary_prompt = type_prompts.get(tipo_resumen, type_prompts["ejecutivo"])

        prompt = (
            f"{summary_prompt}\n\n"
            f"Idioma del resumen: {lang_name}\n"
            f"Documento: {filename}\n\n"
            f"Contenido:\n{content}"
        )

        messages = [{"role": "user", "content": prompt}]
        response = ollama.generate_chat(messages)

        return str(response).strip() if response else ""

    except Exception as e:
        logger.debug(f"Error generando resumen con LLM: {e}")
        return ""


# ============================================================
# S5.4: FORM FILLER / DATA EXTRACTOR
# ============================================================

def extraer_datos(
    ruta: str,
    tipo_documento: str = "auto",
    campos: str = "",
    formato_salida: str = "json"
) -> str:
    """S5.4: Extrae datos estructurados de documentos no estructurados.

    Permite extraer informacion especifica de facturas, CVs,
    contratos, emails u otros documentos, generando una salida
    estructurada en JSON o tabla.

    Args:
        ruta: Ruta al archivo de documento (PDF, DOCX, TXT, CSV)
        tipo_documento: Tipo de documento: auto, factura, cv, contrato, email, general
        campos: Lista de campos a extraer separados por coma (ej: "nombre,fecha,total").
                Si vacio, extrae todos los campos detectables.
        formato_salida: Formato de salida: json o tabla

    Returns:
        Datos extraidos en formato estructurado
    """
    validation = validate_path(ruta)
    if validation != ruta:
        return validation

    # Leer el documento
    ext = os.path.splitext(ruta)[1].lower()
    try:
        if ext == ".pdf":
            content = leer_pdf(ruta)
        elif ext == ".docx":
            content = leer_docx(ruta)
        elif ext in (".txt", ".md"):
            with open(ruta, "r", encoding="utf-8", errors="replace") as f:
                content = f.read()
        elif ext == ".csv":
            content = leer_csv(ruta)
        else:
            content = leer_documento(ruta)
    except Exception as e:
        return f"ERROR: No se pudo leer el archivo: {e}"

    if not content or len(content.strip()) < 10:
        return "ERROR: El documento esta vacio o no contiene texto legible."

    # Truncar si es muy largo
    if len(content) > 12000:
        content = content[:10000] + "\n... [contenido truncado]"

    # Determinar tipo de documento si es auto
    if tipo_documento == "auto":
        content_lower = content.lower()[:2000]
        if any(w in content_lower for w in ["factura", "invoice", "nif", "cif", "importe total", "subtotal"]):
            tipo_documento = "factura"
        elif any(w in content_lower for w in ["curriculum", "experiencia", "educacion", "habilidades", "cv"]):
            tipo_documento = "cv"
        elif any(w in content_lower for w in ["contrato", "clausula", "parte", "jurisdiccion"]):
            tipo_documento = "contrato"
        elif any(w in content_lower for w in ["from:", "to:", "subject:", "dear", "atentamente"]):
            tipo_documento = "email"
        else:
            tipo_documento = "general"

    # Construir prompt de extraccion
    tipo_instrucciones = {
        "factura": (
            "Extrae los siguientes datos de esta factura: "
            "numero_factura, fecha, emisor (nombre y NIF/CIF), "
            "receptor (nombre y NIF/CIF), concepto, subtotal, "
            "impuestos (IVA), total, moneda, fecha_vencimiento. "
            "Si un campo no se encuentra, pon null."
        ),
        "cv": (
            "Extrae los siguientes datos de este curriculum: "
            "nombre_completo, email, telefono, ubicacion, "
            "resumen_profesional, educacion (lista), "
            "experiencia_laboral (lista con empresa, puesto, periodo), "
            "habilidades (lista), idiomas (lista). "
            "Si un campo no se encuentra, pon null."
        ),
        "contrato": (
            "Extrae los siguientes datos de este contrato: "
            "tipo_contrato, partes_involucradas, fecha_inicio, "
            "fecha_fin, objeto_del_contrato, clausulas_principales (lista), "
            "penalidades, jurisdiccion, valor. "
            "Si un campo no se encuentra, pon null."
        ),
        "email": (
            "Extrae los siguientes datos de este email: "
            "remitente, destinatario, fecha, asunto, "
            "accion_solicitada, urgencia (alta/media/baja), "
            "puntos_principales (lista), fecha_limite. "
            "Si un campo no se encuentra, pon null."
        ),
        "general": "Extrae todos los datos clave del documento en formato estructurado.",
    }

    instruccion = tipo_instrucciones.get(tipo_documento, tipo_instrucciones["general"])

    if campos:
        instruccion += f"\nCampos especificos solicitados: {campos}"

    formato_instruccion = (
        "Retorna SOLO un objeto JSON valido, sin markdown ni explicaciones. "
        "Cada campo debe ser una clave del JSON."
        if formato_salida == "json"
        else "Retorna los datos en formato tabla con cabeceras y filas."
    )

    prompt = (
        f"{instruccion}\n\n"
        f"{formato_instruccion}\n\n"
        f"TIPO DE DOCUMENTO: {tipo_documento}\n"
        f"CONTENIDO:\n{content}"
    )

    try:
        from llm import ollama
        messages = [{"role": "user", "content": prompt}]
        response = ollama.generate_chat(messages)
        result = str(response).strip()

        if not result:
            return "ERROR: No se pudieron extraer datos del documento."

        # Limpiar markdown si el LLM lo agrego
        if result.startswith("```"):
            result = result.split("\n", 1)[-1].rsplit("```", 1)[0].strip()

        return result

    except Exception as e:
        logger.debug(f"Error extrayendo datos con LLM: {e}")
        return f"ERROR: No se pudo procesar el documento con IA: {e}"


# ============================================================
# S5.6: LOCAL KNOWLEDGE BASE
# ============================================================

# Directorio para la base de conocimiento
_KB_DIR = os.path.join(REPOS_DIR, "knowledge_base")

def guardar_conocimiento(
    titulo: str,
    contenido: str,
    etiquetas: str = "",
    categoria: str = "general"
) -> str:
    """S5.6: Guarda un fragmento de conocimiento en la base de conocimiento local.

    Almacena notas, fragmentos de documentos importantes, o cualquier
    informacion que el usuario quiera recuperar despues por busqueda
    semantica o por etiquetas.

    Args:
        titulo: Titulo descriptivo del conocimiento
        contenido: Texto del conocimiento a guardar
        etiquetas: Etiquetas separadas por coma (ej: "python,algoritmos,busqueda")
        categoria: Categoria del conocimiento: general, codigo, proyecto, referencia, aprendizaje

    Returns:
        Confirmacion con ID del conocimiento guardado
    """
    import json as _json
    from datetime import datetime as _dt

    # Asegurar directorio
    os.makedirs(_KB_DIR, exist_ok=True)

    # Generar ID unico
    entry_id = _dt.now().strftime("%Y%m%d_%H%M%S")
    entry_file = os.path.join(_KB_DIR, f"{entry_id}.json")

    # Parsear etiquetas
    tag_list = [t.strip().lower() for t in etiquetas.split(",") if t.strip()] if etiquetas else []

    entry = {
        "id": entry_id,
        "titulo": titulo[:200],
        "contenido": contenido[:5000],
        "etiquetas": tag_list,
        "categoria": categoria,
        "timestamp": _dt.now().isoformat(),
        "accesos": 0,
    }

    try:
        with open(entry_file, "w", encoding="utf-8") as f:
            _json.dump(entry, f, ensure_ascii=False, indent=2)

        # Tambien guardar en la memoria de largo plazo del agente
        try:
            from memory.triple_memory import learning
            learning.add_correction(
                f"Conocimiento: {titulo}",
                contenido[:300],
                source="knowledge_base"
            )
        except Exception:
            pass  # No bloquear si la memoria no esta disponible

        return (
            f"Conocimiento guardado exitosamente.\n"
            f"ID: {entry_id}\n"
            f"Titulo: {titulo}\n"
            f"Etiquetas: {', '.join(tag_list) if tag_list else '(sin etiquetas)'}\n"
            f"Categoria: {categoria}"
        )

    except Exception as e:
        return f"ERROR: No se pudo guardar el conocimiento: {e}"


def buscar_conocimiento(
    consulta: str,
    categoria: str = "",
    etiqueta: str = "",
    limite: int = 5
) -> str:
    """S5.6: Busca en la base de conocimiento local por texto, etiqueta o categoria.

    Realiza busqueda semantica en los conocimientos guardados,
    priorizando coincidencias por titulo, etiquetas y contenido.

    Args:
        consulta: Texto a buscar en la base de conocimiento
        categoria: Filtrar por categoria (general, codigo, proyecto, referencia, aprendizaje)
        etiqueta: Filtrar por etiqueta especifica
        limite: Maximo de resultados a retornar

    Returns:
        Lista de conocimientos encontrados con su contenido
    """
    import json as _json

    if not os.path.exists(_KB_DIR):
        return "No hay conocimientos guardados aun. Usa guardar_conocimiento para agregar entradas."

    consulta_lower = consulta.lower()
    entries = []

    try:
        for fname in os.listdir(_KB_DIR):
            if not fname.endswith(".json"):
                continue
            fpath = os.path.join(_KB_DIR, fname)
            try:
                with open(fpath, "r", encoding="utf-8") as f:
                    entry = _json.load(f)
            except Exception:
                continue

            # Filtrar por categoria
            if categoria and entry.get("categoria", "") != categoria:
                continue

            # Filtrar por etiqueta
            if etiqueta:
                entry_tags = [t.lower() for t in entry.get("etiquetas", [])]
                if etiqueta.lower() not in entry_tags:
                    continue

            # Calcular score de relevancia
            score = 0
            titulo = entry.get("titulo", "").lower()
            contenido = entry.get("contenido", "").lower()
            entry_tags = [t.lower() for t in entry.get("etiquetas", [])]

            # Coincidencia en titulo (peso mayor)
            for word in consulta_lower.split():
                if word in titulo:
                    score += 3
                if word in contenido:
                    score += 1
                if word in entry_tags:
                    score += 2

            if score > 0:
                entry["_score"] = score
                entries.append(entry)

    except Exception as e:
        return f"ERROR: Error buscando en base de conocimiento: {e}"

    if not entries:
        return f"No se encontraron conocimientos para: '{consulta}'"

    # Ordenar por score
    entries.sort(key=lambda x: x.get("_score", 0), reverse=True)
    results = entries[:limite]

    # Formatear resultados
    output_lines = [f"Encontrados {len(entries)} resultado(s), mostrando {len(results)}:\n"]
    for i, entry in enumerate(results, 1):
        output_lines.append(f"--- Resultado {i} (relevancia: {entry.get('_score', 0)}) ---")
        output_lines.append(f"ID: {entry.get('id', '?')}")
        output_lines.append(f"Titulo: {entry.get('titulo', '?')}")
        output_lines.append(f"Categoria: {entry.get('categoria', '?')}")
        output_lines.append(f"Etiquetas: {', '.join(entry.get('etiquetas', []))}")
        output_lines.append(f"Contenido: {entry.get('contenido', '')[:500]}")
        output_lines.append(f"Fecha: {entry.get('timestamp', '?')}")
        output_lines.append("")

    return "\n".join(output_lines)


def listar_conocimiento(
    categoria: str = "",
    etiqueta: str = "",
    limite: int = 20
) -> str:
    """S5.6: Lista todos los conocimientos guardados en la base de conocimiento.

    Args:
        categoria: Filtrar por categoria (opcional)
        etiqueta: Filtrar por etiqueta (opcional)
        limite: Maximo de entradas a listar

    Returns:
        Lista de titulos y metadatos de los conocimientos guardados
    """
    import json as _json

    if not os.path.exists(_KB_DIR):
        return "No hay conocimientos guardados aun."

    entries = []
    try:
        for fname in sorted(os.listdir(_KB_DIR), reverse=True):
            if not fname.endswith(".json"):
                continue
            fpath = os.path.join(_KB_DIR, fname)
            try:
                with open(fpath, "r", encoding="utf-8") as f:
                    entry = _json.load(f)
            except Exception:
                continue

            if categoria and entry.get("categoria", "") != categoria:
                continue
            if etiqueta:
                entry_tags = [t.lower() for t in entry.get("etiquetas", [])]
                if etiqueta.lower() not in entry_tags:
                    continue

            entries.append(entry)

    except Exception as e:
        return f"ERROR: Error listando base de conocimiento: {e}"

    if not entries:
        return "No hay conocimientos guardados con esos filtros."

    entries = entries[:limite]
    output_lines = [f"Base de conocimiento: {len(entries)} entrada(s)\n"]
    for entry in entries:
        tags = ", ".join(entry.get("etiquetas", []))
        output_lines.append(
            f"- [{entry.get('id', '?')}] {entry.get('titulo', '?')} "
            f"({entry.get('categoria', '?')}) "
            f"{f'[tags: {tags}]' if tags else ''}"
        )

    return "\n".join(output_lines)
